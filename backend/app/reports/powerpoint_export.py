from __future__ import annotations

import io
import re
import zipfile
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.datasets.storage import DatasetRecord
from app.reports.report_content import collect_report_content, combined_limitations, report_filename, safe_text

_NAVY = RGBColor(23, 37, 84)
_BLUE = RGBColor(37, 99, 235)
_SLATE = RGBColor(71, 85, 105)
_WHITE = RGBColor(255, 255, 255)


def build_powerpoint_report(record: DatasetRecord) -> bytes:
    """Build a bounded executive PowerPoint deck from deterministic evidence."""
    content = collect_report_content(record)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _title_slide(presentation, content.report)
    _overview_slide(presentation, content.report)
    _findings_slide(presentation, content.report)
    _chart_slide(presentation, "Anomalies", content.report.get("anomalies", [])[:10], "column", "anomaly_count")
    _correlation_slide(presentation, content.report.get("correlations", [])[:10])
    _analysis_slide(presentation, content.analysis)
    _limitations_slide(presentation, combined_limitations(content), content)
    _provenance_slide(presentation, content.analysis, content.report)
    output = io.BytesIO()
    presentation.save(output)
    return _normalize_chart_axis_ids(output.getvalue())


def powerpoint_report_filename(record: DatasetRecord) -> str:
    return report_filename(record, "pptx")


def _title_slide(presentation: Presentation, report: dict[str, Any]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _background(slide, _NAVY)
    _text(slide, 0.85, 1.25, 11.6, 1.0, "AI Data Analyst", 50, _WHITE, bold=True)
    _text(slide, 0.85, 2.25, 11.6, 0.7, "Deterministic analysis report", 24, RGBColor(147, 197, 253))
    _text(slide, 0.85, 4.85, 11.6, 0.55, safe_text(report["filename"]), 18, _WHITE, bold=True)
    _text(slide, 0.85, 5.45, 11.6, 0.45, f"Generated {safe_text(report['generated_at'])}", 14, RGBColor(203, 213, 225))


def _overview_slide(presentation: Presentation, report: dict[str, Any]) -> None:
    slide = _content_slide(presentation, "Executive overview")
    overview = report["overview"]
    metrics = [
        ("Rows", f"{overview['rows']:,}"), ("Columns", f"{overview['columns']:,}"),
        ("Missing values", f"{overview['missing_total']:,}"), ("Duplicate rows", f"{overview['duplicate_rows']:,}"),
    ]
    for index, (label, value) in enumerate(metrics):
        x = 0.8 + (index % 2) * 6.1
        y = 1.65 + (index // 2) * 2.25
        _text(slide, x, y, 5.55, 1.5, value, 34, _NAVY, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, x, y + 0.75, 5.55, 0.5, label, 18, _SLATE, align=PP_ALIGN.CENTER)


def _findings_slide(presentation: Presentation, report: dict[str, Any]) -> None:
    slide = _content_slide(presentation, "Key findings")
    findings = report.get("key_findings", [])[:7] or ["No key findings are available."]
    _bullets(slide, findings, top=1.45, font_size=18)


def _chart_slide(presentation: Presentation, title: str, items: list[dict[str, Any]], category: str, value: str) -> None:
    slide = _content_slide(presentation, title)
    if not items:
        _text(slide, 0.9, 2.4, 11.5, 0.7, "No chartable anomaly results are available.", 20, _SLATE)
        return
    data = CategoryChartData()
    data.categories = [safe_text(item.get(category), 50) for item in items]
    data.add_series("Outliers", [float(item.get(value) or 0) for item in items])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(1.45), Inches(11.5), Inches(5.25), data).chart
    _style_chart(chart, "Outlier count")


def _correlation_slide(presentation: Presentation, items: list[dict[str, Any]]) -> None:
    slide = _content_slide(presentation, "Notable correlations")
    if not items:
        _text(slide, 0.9, 2.4, 11.5, 0.7, "No chartable correlation results are available.", 20, _SLATE)
        return
    data = CategoryChartData()
    data.categories = [f"{safe_text(item.get('column_a'), 22)} / {safe_text(item.get('column_b'), 22)}" for item in items]
    data.add_series("Absolute correlation", [abs(float(item.get("correlation") or 0)) for item in items])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.45), Inches(8.3), Inches(5.25), data).chart
    _style_chart(chart, "Absolute r")
    _text(slide, 9.35, 1.8, 3.1, 2.0, "Correlation is descriptive evidence. It does not establish causation.", 18, _NAVY, bold=True)


def _analysis_slide(presentation: Presentation, analysis: Any | None) -> None:
    slide = _content_slide(presentation, "Latest grounded analysis")
    text = "No Deep Reasoning analysis is available in the current server session." if analysis is None else safe_text(analysis.final_answer_text, 2_000)
    _text(slide, 0.9, 1.45, 11.5, 5.25, text, 16, _NAVY)


def _limitations_slide(presentation: Presentation, limitations: list[Any], content: Any) -> None:
    slide = _content_slide(presentation, "Caveats and limitations")
    items = [f"{item.severity}: {item.text}" for item in limitations[:7]]
    if content.caveats.duplicate_row_count:
        items.append(f"Duplicate rows: {content.caveats.duplicate_row_count} ({content.caveats.duplicate_pct}%)")
    items.extend(notice.message for notice in content.ingestion_notices[:2])
    _bullets(slide, items[:9] or ["No limitations were identified for this report."], top=1.4, font_size=16)


def _provenance_slide(presentation: Presentation, analysis: Any | None, report: dict[str, Any]) -> None:
    slide = _content_slide(presentation, "Evidence provenance")
    facts = list(analysis.fact_ledger.facts) if analysis is not None and analysis.fact_ledger else []
    items = [f"{fact.id}: {fact.tool} → {fact.result_path}" for fact in facts[:8]]
    items.extend([f"Dataset ID: {report['dataset_id']}", f"Source file: {safe_text(report['filename'])}"])
    _bullets(slide, items or ["No grounded facts are available."], top=1.4, font_size=16)


def _content_slide(presentation: Presentation, title: str):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _background(slide, _WHITE)
    _text(slide, 0.75, 0.35, 11.9, 0.75, title, 35, _NAVY, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.75), Inches(1.17), Inches(1.3), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = _BLUE; line.line.fill.background()
    _text(slide, 11.7, 6.95, 0.8, 0.25, str(len(presentation.slides)), 10, _SLATE, align=PP_ALIGN.RIGHT)
    return slide


def _bullets(slide: Any, items: list[str], top: float, font_size: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.4), Inches(5.6))
    frame = box.text_frame
    frame.clear(); frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = safe_text(item, 450)
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _NAVY
        paragraph.space_after = Pt(12)
        paragraph.text = "• " + paragraph.text


def _style_chart(chart: Any, axis_title: str) -> None:
    chart.has_legend = False
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = axis_title
    chart.value_axis.maximum_scale = None
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.number_format = "0.0"
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = _BLUE


def _background(slide: Any, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def _normalize_chart_axis_ids(payload: bytes) -> bytes:
    """Rewrite signed python-pptx axis IDs as schema-valid unsigned integers."""
    source = io.BytesIO(payload)
    target = io.BytesIO()
    pattern = re.compile(rb'(<c:(?:axId|crossAx) val=")(-\d+)("/>)')
    with zipfile.ZipFile(source, "r") as archive, zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as output:
        for info in archive.infolist():
            data = archive.read(info.filename)
            if info.filename.startswith("ppt/charts/chart") and info.filename.endswith(".xml"):
                data = pattern.sub(lambda match: match.group(1) + str(int(match.group(2)) % (2**32)).encode() + match.group(3), data)
            output.writestr(info, data)
    return target.getvalue()


def _text(slide: Any, left: float, top: float, width: float, height: float, text: str, size: int, color: RGBColor, bold: bool = False, align: Any = PP_ALIGN.LEFT) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear(); frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = safe_text(text)
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
