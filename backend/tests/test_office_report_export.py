from __future__ import annotations

import io
import re
import zipfile
from datetime import datetime, timezone

import pandas as pd
from docx import Document
from pptx import Presentation

from app.datasets.storage import DatasetRecord
from app.reports.powerpoint_export import build_powerpoint_report, powerpoint_report_filename
from app.reports.word_export import build_word_report, word_report_filename


def _record() -> DatasetRecord:
    return DatasetRecord(
        id="report-dataset",
        original_filename="quarterly results.csv",
        extension=".csv",
        uploaded_at=datetime.now(timezone.utc),
        df=pd.DataFrame(
            {
                "region": ["North", "South", "North", "South"],
                "revenue": [100.0, 120.0, 10_000.0, 110.0],
                "profit": [20.0, 22.0, 2_000.0, None],
            }
        ),
        stored_path="report-dataset.csv",
    )


def test_word_report_has_expected_sections_and_valid_package() -> None:
    payload = build_word_report(_record())
    document = Document(io.BytesIO(payload))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    assert payload.startswith(b"PK")
    assert "Executive overview" in text
    assert "Key findings" in text
    assert "Latest grounded analysis" in text
    assert "Data caveats" in text
    assert "Fact ledger" in text
    assert "Limitations" in text
    assert len(document.tables) >= 4
    assert all(sum(cell.width.inches for cell in table.rows[0].cells) == 6.5 for table in document.tables)
    with zipfile.ZipFile(io.BytesIO(payload)) as package:
        assert "word/document.xml" in package.namelist()


def test_powerpoint_report_is_bounded_and_contains_core_sections() -> None:
    payload = build_powerpoint_report(_record())
    presentation = Presentation(io.BytesIO(payload))
    all_text = "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text_frame")
    )
    assert payload.startswith(b"PK")
    assert 7 <= len(presentation.slides) <= 9
    assert "Executive overview" in all_text
    assert "Key findings" in all_text
    assert "Caveats and limitations" in all_text
    assert "Evidence provenance" in all_text
    with zipfile.ZipFile(io.BytesIO(payload)) as package:
        chart_xml = b"".join(package.read(name) for name in package.namelist() if name.startswith("ppt/charts/chart"))
    assert not re.search(rb'<c:(?:axId|crossAx) val="-', chart_xml)


def test_office_report_filenames_are_safe_ascii() -> None:
    record = _record()
    record.original_filename = "=résumé 2026.csv"
    assert word_report_filename(record) == "_r_sum__2026-analysis-report.docx"
    assert powerpoint_report_filename(record) == "_r_sum__2026-analysis-report.pptx"
    assert word_report_filename(record).isascii()


def test_office_exports_remove_invalid_xml_control_characters() -> None:
    record = _record()
    record.df.columns = ["region\x01", "revenue", "profit"]
    Document(io.BytesIO(build_word_report(record)))
    Presentation(io.BytesIO(build_powerpoint_report(record)))


def test_office_report_endpoints_return_downloadable_files(client, sample_df) -> None:
    buffer = io.BytesIO()
    sample_df.to_csv(buffer, index=False)
    uploaded = client.post("/api/datasets/upload", files={"file": ("sample.csv", buffer.getvalue(), "text/csv")}).json()
    word = client.get(f"/api/reports/{uploaded['dataset_id']}/word")
    powerpoint = client.get(f"/api/reports/{uploaded['dataset_id']}/powerpoint")
    assert word.status_code == 200
    assert word.headers["content-type"].startswith("application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    Document(io.BytesIO(word.content))
    assert powerpoint.status_code == 200
    assert powerpoint.headers["content-type"].startswith("application/vnd.openxmlformats-officedocument.presentationml.presentation")
    Presentation(io.BytesIO(powerpoint.content))


def test_unknown_dataset_office_exports_return_404(client) -> None:
    for endpoint in ("word", "powerpoint"):
        response = client.get(f"/api/reports/missing-dataset/{endpoint}")
        assert response.status_code == 404
        assert response.json()["detail"] == "Dataset 'missing-dataset' not found."
